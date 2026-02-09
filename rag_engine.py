#!/usr/bin/env python3
"""
=============================================================
GDPRag Engine — Mistral Cloud + ChromaDB Locale
=============================================================
Motore RAG GDPR-compliant per documenti aziendali.
Usabile via CLI o importabile dalla web UI (Gradio).

Supporta: PDF, DOCX, DOC, XLSX, XLS, PPTX, HTML, TXT, MD,
          CSV, JSON, ODT, RTF

Autore: Mediaform s.c.r.l.
=============================================================
"""

import os
import sys
import json
import hashlib
import logging
from pathlib import Path
from datetime import datetime
from dataclasses import dataclass, field

from mistralai import Mistral
import chromadb

# ── Parser opzionali (graceful degradation) ──────────────
try:
    import fitz  # PyMuPDF
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    import docx  # python-docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from bs4 import BeautifulSoup
    HAS_HTML = True
except ImportError:
    HAS_HTML = False

try:
    from odf import text as odf_text
    from odf.opendocument import load as odf_load
    HAS_ODF = True
except ImportError:
    HAS_ODF = False

try:
    from striprtf.striprtf import rtf_to_text
    HAS_RTF = True
except ImportError:
    HAS_RTF = False

try:
    import xlrd
    HAS_XLS = True
except ImportError:
    HAS_XLS = False


# ── Logging ──────────────────────────────────────────────
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%H:%M:%S"
)
log = logging.getLogger("gdprag")


# ── Configurazione ───────────────────────────────────────
@dataclass
class RAGConfig:
    """Configurazione completa del motore RAG."""
    # API
    api_key: str = ""
    embed_model: str = "mistral-embed"
    chat_model: str = "mistral-small-latest"

    # ChromaDB
    chroma_path: str = "./chroma_db"
    collection_name: str = "gdprag_docs"

    # Chunking
    chunk_size: int = 1000
    chunk_overlap: int = 200

    # Retrieval
    top_k: int = 5

    # System prompt per la generazione
    system_prompt: str = """Sei GDPRag, un assistente AI aziendale GDPR-compliant.
Rispondi alle domande basandoti ESCLUSIVAMENTE sul contesto fornito dai documenti aziendali.

Regole:
- Rispondi in italiano
- Se l'informazione non è nel contesto, dillo chiaramente
- Cita la fonte (nome file) quando possibile
- Sii preciso e conciso

⚠️ Risposta generata da sistema AI — verificare le informazioni critiche."""


# ── Formati supportati ───────────────────────────────────

SUPPORTED_EXTENSIONS = {
    # Testo puro
    ".txt", ".md", ".csv", ".json", ".log", ".yml", ".yaml", ".xml",
    # Office
    ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx",
    # Web
    ".html", ".htm",
    # Altri
    ".odt", ".rtf",
}


def get_supported_formats_status() -> dict:
    """Ritorna lo stato dei parser disponibili."""
    return {
        "txt/md/csv/json": ("✅", "Sempre disponibile"),
        "pdf":  ("✅" if HAS_PDF  else "❌", "pip install pymupdf"),
        "docx": ("✅" if HAS_DOCX else "❌", "pip install python-docx"),
        "xlsx": ("✅" if HAS_XLSX else "❌", "pip install openpyxl"),
        "xls":  ("✅" if HAS_XLS  else "❌", "pip install xlrd"),
        "pptx": ("✅" if HAS_PPTX else "❌", "pip install python-pptx"),
        "html": ("✅" if HAS_HTML else "❌", "pip install beautifulsoup4"),
        "odt":  ("✅" if HAS_ODF  else "❌", "pip install odfpy"),
        "rtf":  ("✅" if HAS_RTF  else "❌", "pip install striprtf"),
    }


# ── Estrattori di testo ─────────────────────────────────

def extract_text(filepath: Path) -> str:
    """Estrae testo da qualsiasi formato supportato."""
    ext = filepath.suffix.lower()

    try:
        # ── Testo puro ──
        if ext in (".txt", ".md", ".csv", ".log", ".yml", ".yaml", ".xml"):
            return filepath.read_text(encoding="utf-8", errors="ignore")

        if ext == ".json":
            data = json.loads(filepath.read_text(encoding="utf-8"))
            return json.dumps(data, ensure_ascii=False, indent=2)

        # ── PDF ──
        if ext == ".pdf":
            if not HAS_PDF:
                log.warning(f"Skipping {filepath.name} — installa pymupdf")
                return ""
            doc = fitz.open(str(filepath))
            text = "\n\n".join(page.get_text() for page in doc)
            doc.close()
            return text

        # ── Word DOCX ──
        if ext == ".docx":
            if not HAS_DOCX:
                log.warning(f"Skipping {filepath.name} — installa python-docx")
                return ""
            doc = docx.Document(str(filepath))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            # Estrai anche testo dalle tabelle
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        paragraphs.append(row_text)
            return "\n\n".join(paragraphs)

        # ── Word DOC (vecchio formato binario) ──
        if ext == ".doc":
            # Prova con antiword o textract se disponibile
            import subprocess
            try:
                result = subprocess.run(
                    ["antiword", str(filepath)],
                    capture_output=True, text=True, timeout=30
                )
                if result.returncode == 0:
                    return result.stdout
            except FileNotFoundError:
                pass
            log.warning(
                f"Skipping {filepath.name} — file .doc richiede 'antiword' "
                f"(apt install antiword) oppure converti in .docx"
            )
            return ""

        # ── Excel XLSX ──
        if ext == ".xlsx":
            if not HAS_XLSX:
                log.warning(f"Skipping {filepath.name} — installa openpyxl")
                return ""
            wb = openpyxl.load_workbook(str(filepath), read_only=True, data_only=True)
            parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                parts.append(f"--- Foglio: {sheet_name} ---")
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join(
                        str(cell) for cell in row if cell is not None
                    )
                    if row_text.strip():
                        parts.append(row_text)
            wb.close()
            return "\n".join(parts)

        # ── Excel XLS (vecchio formato) ──
        if ext == ".xls":
            if not HAS_XLS:
                log.warning(f"Skipping {filepath.name} — installa xlrd")
                return ""
            wb = xlrd.open_workbook(str(filepath))
            parts = []
            for sheet in wb.sheets():
                parts.append(f"--- Foglio: {sheet.name} ---")
                for row_idx in range(sheet.nrows):
                    row_text = " | ".join(
                        str(cell.value) for cell in sheet.row(row_idx)
                        if cell.value
                    )
                    if row_text.strip():
                        parts.append(row_text)
            return "\n".join(parts)

        # ── PowerPoint PPTX ──
        if ext == ".pptx":
            if not HAS_PPTX:
                log.warning(f"Skipping {filepath.name} — installa python-pptx")
                return ""
            prs = Presentation(str(filepath))
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    parts.append(f"--- Slide {i} ---")
                    parts.extend(slide_texts)
            return "\n\n".join(parts)

        # ── HTML ──
        if ext in (".html", ".htm"):
            if HAS_HTML:
                html = filepath.read_text(encoding="utf-8", errors="ignore")
                soup = BeautifulSoup(html, "html.parser")
                # Rimuovi script e style
                for tag in soup(["script", "style"]):
                    tag.decompose()
                return soup.get_text(separator="\n", strip=True)
            else:
                # Fallback: leggi come testo grezzo
                return filepath.read_text(encoding="utf-8", errors="ignore")

        # ── ODT (LibreOffice/OpenOffice) ──
        if ext == ".odt":
            if not HAS_ODF:
                log.warning(f"Skipping {filepath.name} — installa odfpy")
                return ""
            doc = odf_load(str(filepath))
            paragraphs = doc.getElementsByType(odf_text.P)
            texts = []
            for p in paragraphs:
                t = ""
                for node in p.childNodes:
                    if hasattr(node, "data"):
                        t += node.data
                    elif hasattr(node, "__str__"):
                        t += str(node)
                if t.strip():
                    texts.append(t.strip())
            return "\n\n".join(texts)

        # ── RTF ──
        if ext == ".rtf":
            if not HAS_RTF:
                log.warning(f"Skipping {filepath.name} — installa striprtf")
                return ""
            raw = filepath.read_text(encoding="utf-8", errors="ignore")
            return rtf_to_text(raw)

        log.warning(f"Formato non supportato: {filepath.name}")
        return ""

    except Exception as e:
        log.error(f"Errore leggendo {filepath.name}: {e}")
        return ""


# ── Motore RAG ───────────────────────────────────────────

class RAGEngine:
    """Motore RAG completo: ingest, query, gestione documenti."""

    def __init__(self, config: RAGConfig = None):
        self.config = config or RAGConfig()
        self.client = None
        self.chroma = None
        self.collection = None

    def _get_client(self) -> Mistral:
        """Lazy init del client Mistral."""
        if self.client is None:
            api_key = self.config.api_key or os.environ.get("MISTRAL_API_KEY", "")
            if not api_key:
                raise ValueError(
                    "API key mancante. Imposta MISTRAL_API_KEY o passa api_key nella config."
                )
            self.client = Mistral(api_key=api_key)
        return self.client

    def _get_chroma(self) -> chromadb.ClientAPI:
        """Lazy init di ChromaDB."""
        if self.chroma is None:
            self.chroma = chromadb.PersistentClient(path=self.config.chroma_path)
        return self.chroma

    def _get_collection(self, create: bool = False):
        """Ottieni o crea la collection."""
        chroma = self._get_chroma()
        if create:
            return chroma.get_or_create_collection(
                name=self.config.collection_name,
                metadata={"hnsw:space": "cosine"}
            )
        return chroma.get_collection(self.config.collection_name)

    # ── Caricamento documenti ──

    def load_sources(self, paths: list[str]) -> list[dict]:
        """
        Carica documenti da una lista di percorsi (file o cartelle).
        Le cartelle vengono scansionate ricorsivamente.
        """
        docs = []
        for path_str in paths:
            path = Path(path_str).expanduser().resolve()

            if not path.exists():
                log.warning(f"Percorso non trovato: {path}")
                continue

            if path.is_file():
                files = [path]
            else:
                files = sorted(path.rglob("*"))

            for file in files:
                if file.is_dir():
                    continue
                if file.suffix.lower() not in SUPPORTED_EXTENSIONS:
                    continue
                if file.name.startswith(".") or file.name.startswith("~"):
                    continue  # Ignora file nascosti e temp di Office

                log.info(f"Caricamento: {file.name}")
                text = extract_text(file)

                if text.strip():
                    # Hash del contenuto per dedup
                    content_hash = hashlib.md5(text.encode()).hexdigest()[:12]
                    docs.append({
                        "filename": file.name,
                        "filepath": str(file),
                        "source_dir": str(file.parent),
                        "text": text.strip(),
                        "hash": content_hash,
                        "size": len(text),
                        "modified": datetime.fromtimestamp(
                            file.stat().st_mtime
                        ).isoformat()
                    })
                    log.info(f"  ✅ {file.name} ({len(text):,} caratteri)")
                else:
                    log.warning(f"  ⚠️  {file.name} — nessun testo estratto")

        log.info(f"Totale: {len(docs)} documenti caricati")
        return docs

    def load_sources_file(self, sources_file: str) -> list[dict]:
        """
        Carica documenti da un file sources.txt.
        Ogni riga è un percorso (file o cartella).
        Righe vuote e commenti (#) vengono ignorate.
        """
        paths = []
        sf = Path(sources_file)
        if not sf.exists():
            raise FileNotFoundError(f"File sorgenti non trovato: {sources_file}")

        for line in sf.read_text().splitlines():
            line = line.strip()
            if line and not line.startswith("#"):
                paths.append(line)

        if not paths:
            raise ValueError(f"Nessun percorso trovato in {sources_file}")

        log.info(f"Sorgenti da {sources_file}: {len(paths)} percorsi")
        return self.load_sources(paths)

    # ── Chunking ──

    def chunk_text(self, text: str) -> list[str]:
        """Divide il testo in chunk con sovrapposizione intelligente."""
        chunks = []
        start = 0
        size = self.config.chunk_size
        overlap = self.config.chunk_overlap

        while start < len(text):
            end = start + size

            if end < len(text):
                for sep in ["\n\n", "\n", ". ", "! ", "? ", "; "]:
                    last_sep = text.rfind(sep, start + size // 2, end)
                    if last_sep != -1:
                        end = last_sep + len(sep)
                        break

            chunk = text[start:end].strip()
            if chunk and len(chunk) > 50:  # Ignora chunk troppo corti
                chunks.append(chunk)

            start = end - overlap
            if start >= len(text):
                break

        return chunks

    def chunk_documents(self, docs: list[dict]) -> tuple[list[str], list[dict]]:
        """Chunk tutti i documenti, ritorna (chunks, metadati)."""
        all_chunks = []
        all_metadata = []

        for doc in docs:
            chunks = self.chunk_text(doc["text"])
            for i, chunk in enumerate(chunks):
                all_chunks.append(chunk)
                all_metadata.append({
                    "filename": doc["filename"],
                    "source_dir": doc["source_dir"],
                    "chunk_index": i,
                    "total_chunks": len(chunks),
                    "content_hash": doc["hash"],
                    "doc_modified": doc["modified"]
                })

        log.info(f"Chunking: {len(all_chunks)} chunk da {len(docs)} documenti")
        return all_chunks, all_metadata

    # ── Embedding ──

    def embed_texts(self, texts: list[str], batch_size: int = 50) -> list[list[float]]:
        """Genera embedding via API Mistral, in batch."""
        client = self._get_client()
        all_embeddings = []
        total = len(texts)

        for i in range(0, total, batch_size):
            batch = texts[i:i + batch_size]
            batch_num = i // batch_size + 1
            total_batches = (total + batch_size - 1) // batch_size
            log.info(f"Embedding batch {batch_num}/{total_batches} ({len(batch)} chunk)")

            response = client.embeddings.create(
                model=self.config.embed_model,
                inputs=batch
            )
            batch_embeddings = [item.embedding for item in response.data]
            all_embeddings.extend(batch_embeddings)

        log.info(f"Embedding completato: {len(all_embeddings)} vettori")
        return all_embeddings

    # ── Ingest ──

    def ingest(self, paths: list[str] = None, sources_file: str = None,
               append: bool = False,
               progress_callback=None) -> dict:
        """
        Pipeline di ingestione completa.

        Args:
            paths: lista di percorsi (file o cartelle)
            sources_file: percorso a un file sources.txt
            append: se True, aggiunge ai documenti esistenti
            progress_callback: funzione(step, message) per UI

        Returns:
            dict con statistiche dell'ingestione
        """
        def progress(step, msg):
            log.info(msg)
            if progress_callback:
                progress_callback(step, msg)

        stats = {
            "documents": 0, "chunks": 0, "tokens_est": 0,
            "cost_est": 0.0, "errors": []
        }

        # 1. Carica documenti
        progress("load", "📁 Caricamento documenti...")
        if sources_file:
            docs = self.load_sources_file(sources_file)
        elif paths:
            docs = self.load_sources(paths)
        else:
            raise ValueError("Specifica paths o sources_file")

        if not docs:
            stats["errors"].append("Nessun documento trovato")
            return stats

        stats["documents"] = len(docs)

        # 2. Chunking
        progress("chunk", "✂️ Chunking documenti...")
        chunks, metadata = self.chunk_documents(docs)
        stats["chunks"] = len(chunks)

        # 3. Embedding
        progress("embed", f"☁️ Embedding {len(chunks)} chunk via Mistral API...")
        embeddings = self.embed_texts(chunks)

        # 4. Storage
        progress("store", "💾 Salvataggio in ChromaDB...")
        chroma = self._get_chroma()

        if not append:
            try:
                chroma.delete_collection(self.config.collection_name)
                log.info("Collection precedente eliminata")
            except Exception:
                pass

        collection = self._get_collection(create=True)

        # Genera ID univoci (basati su hash del contenuto)
        existing_count = collection.count() if append else 0
        ids = [f"chunk_{existing_count + i:06d}" for i in range(len(chunks))]

        # ChromaDB ha un limite di ~41666 per batch
        batch_limit = 5000
        for i in range(0, len(chunks), batch_limit):
            end = min(i + batch_limit, len(chunks))
            collection.add(
                ids=ids[i:end],
                documents=chunks[i:end],
                embeddings=embeddings[i:end],
                metadatas=metadata[i:end]
            )

        # Statistiche
        total_chars = sum(len(c) for c in chunks)
        stats["tokens_est"] = total_chars // 4
        stats["cost_est"] = stats["tokens_est"] * 0.10 / 1_000_000

        total_chunks = collection.count()
        progress("done",
                 f"✅ Completato! {stats['documents']} documenti → "
                 f"{stats['chunks']} chunk nuovi (totale in DB: {total_chunks}). "
                 f"Costo embedding: ${stats['cost_est']:.4f}")

        return stats

    # ── Query ──

    def query(self, question: str, top_k: int = None) -> dict:
        """
        Pipeline RAG: embed domanda → ricerca → generazione.

        Returns:
            dict con: answer, sources, usage
        """
        if top_k is None:
            top_k = self.config.top_k

        client = self._get_client()
        collection = self._get_collection()

        # 1. Embed domanda
        q_response = client.embeddings.create(
            model=self.config.embed_model,
            inputs=[question]
        )
        q_embedding = q_response.data[0].embedding

        # 2. Ricerca in ChromaDB
        results = collection.query(
            query_embeddings=[q_embedding],
            n_results=top_k,
            include=["documents", "metadatas", "distances"]
        )

        # 3. Prepara contesto
        sources = []
        context_parts = []
        for doc, meta, dist in zip(
            results["documents"][0],
            results["metadatas"][0],
            results["distances"][0]
        ):
            similarity = 1 - dist
            source_info = {
                "filename": meta.get("filename", "?"),
                "chunk_index": meta.get("chunk_index", "?"),
                "similarity": round(similarity, 4),
                "preview": doc[:200] + "..." if len(doc) > 200 else doc
            }
            sources.append(source_info)
            context_parts.append(
                f"[Fonte: {meta['filename']}, sezione {meta['chunk_index']}]\n{doc}"
            )

        context = "\n\n---\n\n".join(context_parts)

        # 4. Genera risposta
        user_prompt = f"""Contesto dai documenti aziendali:
---------------------
{context}
---------------------

Domanda: {question}

Rispondi basandoti solo sul contesto fornito sopra."""

        response = client.chat.complete(
            model=self.config.chat_model,
            messages=[
                {"role": "system", "content": self.config.system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.1,
            max_tokens=1024
        )

        answer = response.choices[0].message.content
        usage = {
            "prompt_tokens": response.usage.prompt_tokens,
            "completion_tokens": response.usage.completion_tokens,
            "total_tokens": response.usage.total_tokens
        }

        return {
            "answer": answer,
            "sources": sources,
            "usage": usage
        }

    # ── Statistiche ──

    def get_stats(self) -> dict:
        """Ritorna statistiche sulla collection."""
        try:
            collection = self._get_collection()
            count = collection.count()

            # Calcola dimensione DB
            db_path = Path(self.config.chroma_path)
            if db_path.exists():
                db_size = sum(
                    f.stat().st_size for f in db_path.rglob("*") if f.is_file()
                ) / 1024 / 1024
            else:
                db_size = 0

            # Ottieni lista file unici
            sample = collection.peek(limit=min(count, 100))
            filenames = set()
            if sample.get("metadatas"):
                for meta in sample["metadatas"]:
                    filenames.add(meta.get("filename", "?"))

            return {
                "total_chunks": count,
                "db_size_mb": round(db_size, 1),
                "sample_files": sorted(filenames),
                "collection_name": self.config.collection_name
            }
        except Exception as e:
            return {"error": str(e)}

    def list_indexed_files(self) -> list[str]:
        """Ritorna la lista di tutti i file indicizzati."""
        try:
            collection = self._get_collection()
            count = collection.count()
            if count == 0:
                return []

            # Recupera tutti i metadati (in batch se necessario)
            filenames = set()
            batch_size = 1000
            for offset in range(0, count, batch_size):
                result = collection.get(
                    limit=batch_size,
                    offset=offset,
                    include=["metadatas"]
                )
                for meta in result["metadatas"]:
                    filenames.add(meta.get("filename", "?"))

            return sorted(filenames)
        except Exception:
            return []

    def clear(self):
        """Cancella tutta la collection."""
        chroma = self._get_chroma()
        try:
            chroma.delete_collection(self.config.collection_name)
            log.info("Collection eliminata")
        except Exception:
            pass


# ── CLI ──────────────────────────────────────────────────

def main():
    import argparse

    parser = argparse.ArgumentParser(
        description="GDPRag Engine — Mistral Cloud + ChromaDB Locale",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Esempi:
  # Indicizza una cartella
  %(prog)s --ingest --docs-path /percorso/documenti

  # Indicizza da più cartelle
  %(prog)s --ingest --docs-path /cartella1 /cartella2 /cartella3

  # Indicizza da file sources.txt
  %(prog)s --ingest --sources sources.txt

  # Aggiungi documenti senza cancellare i precedenti
  %(prog)s --ingest --docs-path /nuovi_documenti --append

  # Fai una domanda
  %(prog)s --query "Quali sono le procedure di sicurezza?"

  # Chat interattiva
  %(prog)s --interactive

  # Mostra statistiche
  %(prog)s --stats

  # Mostra formati supportati
  %(prog)s --formats
"""
    )

    parser.add_argument("--ingest", action="store_true",
                        help="Indicizza documenti")
    parser.add_argument("--docs-path", type=str, nargs="+",
                        help="Cartelle/file da indicizzare (puoi passarne più di una)")
    parser.add_argument("--sources", type=str,
                        help="File con lista percorsi (uno per riga)")
    parser.add_argument("--append", action="store_true",
                        help="Aggiungi ai documenti esistenti senza cancellare")
    parser.add_argument("--query", type=str,
                        help="Fai una domanda ai documenti")
    parser.add_argument("--interactive", action="store_true",
                        help="Modalità chat interattiva")
    parser.add_argument("--stats", action="store_true",
                        help="Mostra statistiche")
    parser.add_argument("--formats", action="store_true",
                        help="Mostra formati supportati e stato parser")
    parser.add_argument("--clear", action="store_true",
                        help="Cancella tutti i documenti indicizzati")
    parser.add_argument("--files", action="store_true",
                        help="Elenca i file indicizzati")

    args = parser.parse_args()

    # ── Mostra formati ──
    if args.formats:
        print("\n📋 Formati supportati:\n")
        for fmt, (status, note) in get_supported_formats_status().items():
            print(f"  {status} {fmt:12s} {note}")
        print()
        return

    config = RAGConfig()
    engine = RAGEngine(config)

    # ── Stats (non richiede API key) ──
    if args.stats:
        stats = engine.get_stats()
        if "error" in stats:
            print(f"❌ {stats['error']}")
        else:
            print(f"\n📊 Statistiche Collection: {stats['collection_name']}")
            print(f"   Chunk indicizzati: {stats['total_chunks']}")
            print(f"   Dimensione DB: {stats['db_size_mb']} MB")
            if stats["sample_files"]:
                print(f"   File indicizzati (campione): {', '.join(stats['sample_files'][:20])}")
        return

    if args.files:
        files = engine.list_indexed_files()
        if files:
            print(f"\n📁 File indicizzati ({len(files)}):\n")
            for f in files:
                print(f"  📄 {f}")
        else:
            print("❌ Nessun file indicizzato")
        return

    if args.clear:
        confirm = input("⚠️  Sicuro di voler cancellare tutti i documenti? (s/N): ")
        if confirm.lower() in ("s", "si", "sì", "y", "yes"):
            engine.clear()
            print("✅ Collection eliminata")
        return

    # ── Ingest ──
    if args.ingest:
        paths = args.docs_path or ["./documenti"]
        try:
            stats = engine.ingest(
                paths=paths if not args.sources else None,
                sources_file=args.sources,
                append=args.append
            )
            if stats.get("errors"):
                for err in stats["errors"]:
                    print(f"⚠️  {err}")
        except Exception as e:
            print(f"❌ Errore: {e}")
        return

    # ── Query singola ──
    if args.query:
        try:
            result = engine.query(args.query)
            print(f"\n{'=' * 60}")
            print(f"📖 Risposta:\n\n{result['answer']}")
            print(f"\n{'─' * 60}")
            print("📚 Fonti:")
            for s in result["sources"]:
                print(f"   • {s['filename']} (sezione {s['chunk_index']}, "
                      f"similarità {s['similarity']:.0%})")
            print(f"\n💰 Token: {result['usage']['total_tokens']}")
            print(f"{'=' * 60}")
        except Exception as e:
            print(f"❌ Errore: {e}")
        return

    # ── Chat interattiva ──
    if args.interactive:
        print("\n" + "=" * 60)
        print("💬 CHAT INTERATTIVA — GDPRag")
        print("=" * 60)
        print("Scrivi le tue domande. 'esci' per uscire.\n")

        while True:
            try:
                question = input("📝 > ").strip()
            except (EOFError, KeyboardInterrupt):
                print("\n👋 Arrivederci!")
                break

            if not question:
                continue
            if question.lower() in ("quit", "exit", "esci", "q"):
                print("👋 Arrivederci!")
                break

            try:
                result = engine.query(question)
                print(f"\n{'─' * 60}")
                print(f"{result['answer']}")
                print(f"\n📚 Fonti: {', '.join(s['filename'] for s in result['sources'])}")
                print(f"{'─' * 60}\n")
            except Exception as e:
                print(f"❌ Errore: {e}\n")
        return

    parser.print_help()


if __name__ == "__main__":
    main()
